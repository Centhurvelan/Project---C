import os
import shutil
import zipfile
import pandas as pd
import io
import json
import base64
import mimetypes
import collections
from tenacity import retry, wait_random_exponential, stop_after_attempt, retry_if_exception_type
 
# Document parsing imports
import docx
from pypdf import PdfReader
from pptx import Presentation
# Excel styling imports
import xlsxwriter # Ensure this is installed: pip install XlsxWriter
 
# --- Constants for File Types and AI ---
ALLOWED_IMAGE_EXTENSIONS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp')
ALLOWED_VIDEO_EXTENSIONS = ('.mp4', '.avi', '.mov', '.wmv', '.flv', '.webm')
 
TEXT_FILE_EXTENSIONS = (
    '.py', '.java', '.js', '.ts', '.jsx', '.tsx', # Code
    '.txt', '.md', # Text and Markdown
    '.json', '.yaml', '.yml', '.xml', '.ini', '.cfg', '.conf', # Config/Data
    '.html', '.css', '.sh', # Web and Scripting
    '.env', '.log', # Environment and Logs
    '.c', '.cpp', '.h', '.hpp', # C/C++ family
    '.rb', '.php', '.go', '.cs', '.swift' # Other common languages
)
 
DOCUMENT_PROJECT_EXTENSIONS = ('.docx', '.pdf', '.pptx')
MAX_FILE_SIZE_FOR_AI_PROCESSING = 1000 * 1024 * 1024 # 1 MB
MAX_INDIVIDUAL_FILE_TRUNCATION_CHARS = 4000
MAX_TOTAL_AI_TEXT_CHARS = 200000
 
# --- DYNAMIC RUBRIC KEYWORDS (FINAL & COMPLETE) ---
# These lists are the core of the dynamic handling.
 
POTENTIAL_CRITERION_COLS = [
    'Evaluation Criteria', 'Criterion', 'Item',
    'Code Snippet to be checked / Expected Result',
     'Description'
]
POTENTIAL_PARAMETERS_COLS = [
    'Parameters',
    'WebPage / Class Affected /Test scenario'
]
POTENTIAL_CATEGORY_COLS = [
    'Category', 'Group', 'Module',
    'Skill Cluster', 'Business Requirement'
]
POTENTIAL_MAX_SCORE_COLS = [
    'Score', 'Max Score', 'Points',
    'Weightage', 'Points Possible', 'Max'
]
 
# --- Helper Functions for File Handling and Processing ---
 
def unzip_file(zip_path, extract_to):
    """Unzips a file to a specified directory."""
    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(extract_to)
        return True
    except Exception as e:
        print(f"Error unzipping file {zip_path}: {e}")
        return False
 
def read_docx(file_path):
    """Extracts text from a DOCX file, returns None on error."""
    try:
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
        return None
 
def read_pdf(file_path):
    """Extracts text from a PDF file, returns None on error."""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += (page.extract_text() or "") + "\n"
        return text
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
        return None
 
def read_pptx(file_path):
    """Extracts text from a PPTX file, returns None on error."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPTX {file_path}: {e}")
        return None
    return text
 
def _identify_rubric_columns(df_rubric):
    """
    (REWRITTEN) Dynamically identifies all key columns (Criterion, Parameters, Category, Score)
    by searching for keywords and ensuring no column is used for more than one role.
    """
    standardized_column_map = {}
    used_columns = set()
 
    # Define search order and mapping
    searches = {
        'max_score_col': POTENTIAL_MAX_SCORE_COLS, # Search for numeric score col first
        'category_col': POTENTIAL_CATEGORY_COLS,
        'parameters_col': POTENTIAL_PARAMETERS_COLS,
        'criterion_col': POTENTIAL_CRITERION_COLS
    }
 
    for key, keywords in searches.items():
        if key in standardized_column_map: continue # Already found
        for col in df_rubric.columns:
            if col in used_columns: continue # Already assigned to another role
           
            col_lower = str(col).lower()
            for kw in keywords:
                if kw.lower() in col_lower:
                    # Special check for score to ensure it's a numeric column
                    if key == 'max_score_col':
                        if pd.to_numeric(df_rubric[col], errors='coerce').notna().sum() < len(df_rubric) * 0.5:
                            continue # Not a score column if less than 50% of values are numeric
                   
                    standardized_column_map[key] = col
                    used_columns.add(col)
                    break # Go to next key
            if key in standardized_column_map:
                break # Go to next key

    # Fallback logic if a critical column is not found
    if 'criterion_col' not in standardized_column_map and not df_rubric.empty:
        # Assume the first non-numeric, non-used column is the criterion
        for col in df_rubric.columns:
            if col not in used_columns and not pd.to_numeric(df_rubric[col], errors='coerce').all():
                standardized_column_map['criterion_col'] = col
                print(f"Fallback: Assuming '{col}' is the criterion column.")
                break
   
    if 'criterion_col' not in standardized_column_map:
        print("CRITICAL ERROR: No valid criterion column could be identified.")
        return None
       
    return standardized_column_map
 
def process_rubric_excel(file_path):
    """
    (UPDATED) Processes the rubric Excel/CSV file, dynamically detecting header and ALL key columns.
    """
    try:
        file_ext = os.path.splitext(file_path)[1].lower()
        df_raw = pd.read_csv(file_path, header=None, on_bad_lines='skip') if file_ext == '.csv' else pd.read_excel(file_path, header=None)
       
        # Auto-detect header
        header_keywords = [kw.lower() for kws in [POTENTIAL_CATEGORY_COLS, POTENTIAL_PARAMETERS_COLS, POTENTIAL_CRITERION_COLS, POTENTIAL_MAX_SCORE_COLS] for kw in kws]
        header_row_index = 0
        for i, row in df_raw.head(10).iterrows():
            row_str = ' '.join(row.dropna().astype(str).str.lower().tolist())
            if sum(1 for keyword in header_keywords if keyword in row_str) >= 2: # Find rows with at least 2 keywords
                header_row_index = i
                print(f"Detected header row at index: {header_row_index}")
                break
        else:
            print("Warning: No strong header keywords found. Assuming header is at first row (index 0).")
 
        df_rubric = pd.read_csv(file_path, header=header_row_index) if file_ext == '.csv' else pd.read_excel(file_path, header=header_row_index)
        df_rubric.dropna(axis=1, how='all', inplace=True)
       
        col_map = _identify_rubric_columns(df_rubric)
        if col_map is None: return None, None
        print(f"Identified columns: {col_map}")
 
        # Get dynamically identified column names
        actual_criterion_col = col_map.get('criterion_col')
        actual_category_col = col_map.get('category_col')
        actual_parameters_col = col_map.get('parameters_col')
        actual_max_score_col = col_map.get('max_score_col')
 
        # Forward-fill grouping columns
        grouping_cols = [col for col in [actual_category_col, actual_parameters_col] if col]
        for col in grouping_cols:
            df_rubric[col] = df_rubric[col].ffill()
       
        if actual_max_score_col and grouping_cols:
            df_rubric[actual_max_score_col] = df_rubric.groupby(grouping_cols, sort=False)[actual_max_score_col].ffill()
 
        df_rubric.reset_index(inplace=True)
       
        # Identify and filter out summary rows (where the main criterion is blank)
        df_rubric['is_summary_row'] = df_rubric[actual_criterion_col].isna() | (df_rubric[actual_criterion_col].astype(str).str.strip() == '')
       
        df_gradeable = df_rubric[~df_rubric['is_summary_row']].copy()
        df_for_ai_markdown = df_gradeable.drop(columns=['is_summary_row', 'index'], errors='ignore')
        rubric_string_for_ai = df_for_ai_markdown.to_markdown(index=False)
       
        df_rubric.set_index('index', inplace=True, drop=False)
        df_rubric._identified_columns = col_map
        return rubric_string_for_ai, df_rubric
    except Exception as e:
        print(f"Error processing rubric file: {e}")
        import traceback
        traceback.print_exc()
        return None, None
 
# *** FUNCTION RE-INSERTED HERE TO FIX THE ERROR ***
def encode_image_to_base64(image_path):
    """Encodes an image file to a base64 string for embedding in the prompt."""
    try:
        with open(image_path, 'rb') as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    except Exception as e:
        print(f"Error encoding image {image_path}: {e}")
        return None
 
def collect_project_content(top_level_extracted_base_dir):
    all_text_file_candidates, image_messages_for_ai, video_files_detected = [], [], []
    image_count = 0
    scan_queue = collections.deque([top_level_extracted_base_dir])
    processed_zip_archives = set()
    while scan_queue:
        current_dir_to_scan = scan_queue.popleft()
        for entry in os.scandir(current_dir_to_scan):
            item_path, item_name = entry.path, entry.name
            relative_file_path = os.path.relpath(item_path, top_level_extracted_base_dir)
            if entry.is_dir():
                if item_name not in ['__pycache__', '.idea', '.venv', 'node_modules', '.git', 'dist', 'build']:
                    scan_queue.append(item_path)
            elif "__MACOSX" in item_path or item_name.startswith("._") or item_name == ".DS_Store":
                continue
            elif item_name.lower().endswith('.zip'):
                if os.path.abspath(item_path) not in processed_zip_archives:
                    nested_extract_dir = os.path.join(os.path.dirname(item_path), os.path.splitext(item_name)[0] + "_extracted_nested")
                    os.makedirs(nested_extract_dir, exist_ok=True)
                    if unzip_file(item_path, nested_extract_dir):
                        scan_queue.append(nested_extract_dir)
                    processed_zip_archives.add(os.path.abspath(item_path))
            elif entry.is_file() and entry.stat().st_size > MAX_FILE_SIZE_FOR_AI_PROCESSING:
                print(f"Skipping file (too large): {relative_file_path}")
            elif item_name.lower().endswith(ALLOWED_IMAGE_EXTENSIONS) and image_count < 5:
                encoded_image = encode_image_to_base64(item_path)
                if encoded_image:
                    mime_type = mimetypes.guess_type(item_name)[0] or 'image/jpeg'
                    image_messages_for_ai.append({"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{encoded_image}", "detail": "auto"}})
                    image_count += 1
            elif item_name.lower().endswith(TEXT_FILE_EXTENSIONS + DOCUMENT_PROJECT_EXTENSIONS):
                content = None
                try:
                    if item_name.lower().endswith(TEXT_FILE_EXTENSIONS):
                        with open(item_path, 'r', encoding='utf-8', errors='ignore') as f: content = f.read()
                    elif item_name.lower().endswith('.pdf'): content = read_pdf(item_path)
                    elif item_name.lower().endswith('.docx'): content = read_docx(item_path)
                    elif item_name.lower().endswith('.pptx'): content = read_pptx(item_path)
                    if content: all_text_file_candidates.append({"path": relative_file_path, "content": content})
                except Exception as e:
                    print(f"Error processing file {relative_file_path}: {e}")
            elif item_name.lower().endswith(ALLOWED_VIDEO_EXTENSIONS):
                video_files_detected.append(relative_file_path)
    collected_text_for_ai = {}
    current_total_text_chars = 0
    all_text_file_candidates.sort(key=lambda x: len(x['content']))
    for file_info in all_text_file_candidates:
        truncated_content = file_info['content'][:MAX_INDIVIDUAL_FILE_TRUNCATION_CHARS]
        if current_total_text_chars + len(truncated_content) <= MAX_TOTAL_AI_TEXT_CHARS:
            collected_text_for_ai[file_info['path']] = truncated_content
            current_total_text_chars += len(truncated_content)
        else: break
    return collected_text_for_ai, image_messages_for_ai, video_files_detected
 
def safe_numeric_score(score_input):
    if score_input is None or pd.isna(score_input): return 0.0
    try:
        if isinstance(score_input, str):
            if '/' in score_input: score_input = score_input.split('/')[0].strip()
            return float(score_input)
        return float(score_input)
    except (ValueError, TypeError): return 0.0
 
@retry(wait=wait_random_exponential(multiplier=1, min=4, max=60), stop=stop_after_attempt(5))
def _call_openai_with_retries(chat_client, messages, model_name, temperature, max_tokens, response_format):
    import openai
    try:
        response = chat_client.chat.completions.create(model=model_name, messages=messages, temperature=temperature, max_tokens=max_tokens, response_format=response_format)
        return response
    except (openai.APIConnectionError, openai.RateLimitError, openai.APITimeoutError) as e:
        print(f"OpenAI API error (retriable): {e}"); raise
    except Exception as e:
        print(f"Unexpected error during OpenAI call: {e}"); raise
 
def generate_grading_with_openai(chat_client, original_rubric_dataframe, rubric_data_markdown_for_ai, requirements_text, project_text_files_content, image_messages_for_ai, has_video):
    if not chat_client: return "Azure OpenAI chat client not initialized.", [], {"total_score": "N/A", "overall_feedback": "AI grading skipped."}
    col_map = getattr(original_rubric_dataframe, '_identified_columns', {}); actual_criteria_col_name = col_map.get('criterion_col'); actual_max_score_col_name = col_map.get('max_score_col')
    if not actual_criteria_col_name: return "Failed to identify grading criteria.", [], {"total_score": "N/A", "overall_feedback": "Could not identify grading criteria."}
    criteria_for_ai_list = []
    gradeable_rubric = original_rubric_dataframe[~original_rubric_dataframe['is_summary_row']].copy()
    gradeable_rubric.dropna(subset=[actual_criteria_col_name], inplace=True)
    for index, row in gradeable_rubric.iterrows():
        criterion_name = str(row.get(actual_criteria_col_name, "")).strip()
        criterion_info = {"criterion_id": index, "criterion_name": criterion_name}
        if actual_max_score_col_name and pd.notna(row.get(actual_max_score_col_name)):
            criterion_info["max_score"] = safe_numeric_score(row[actual_max_score_col_name])
        for col, val in row.items():
            if col not in [actual_criteria_col_name, actual_max_score_col_name, 'is_summary_row', 'index'] and pd.notna(val):
                criterion_info[str(col).replace(" ", "_").lower()] = str(val).strip()
        criteria_for_ai_list.append(criterion_info)
    criteria_list_str = json.dumps(criteria_for_ai_list, indent=2)
    video_guidance_text = "Video file detected. Assume video-related criteria are met." if has_video else ""
    main_prompt_text = f"""You are an expert software project grader. Evaluate a project based on the rubric, requirements, and project files.
**Evaluation Rubric (for context):**\n{rubric_data_markdown_for_ai}
**List of Specific Criteria to Grade:**\nThis is the definitive list. You MUST provide a grade for EACH object in this JSON array.\n```json\n{criteria_list_str}\n```
**Project Requirements:**\n{requirements_text}
**Project Content:**\n(Code, configs, etc. from the project submission follow)\n{''.join(f"File: {filename}\\n```\\n{content}\\n```\\n" for filename, content in project_text_files_content.items())}End of Project Content.
**Visual Analysis:**\n{len(image_messages_for_ai)} UI screenshots are provided. {video_guidance_text}
---
**Grading Task:**\nProvide a grade for EACH criterion from the "List of Specific Criteria to Grade".
**Output your response STRICTLY as a single JSON object. Do not include any other text.**
The JSON object must have this structure:
```json
{{
    "overall_total_score": "Total calculated score (numeric if possible)",
    "overall_feedback": "A comprehensive summary of the project's performance.",
    "grades": [ {{ "criterion_id": 0, "criterion_name": "The EXACT name of the criterion from the list", "score_achieved": 4.0, "comments": "Specific justification for this score." }} ]
}}
```
**CRITICAL:** The `criterion_id` in your output MUST EXACTLY MATCH the `criterion_id` from the list I provided. This is essential for matching results. The `criterion_name` should also be returned exactly as provided.
"""
    messages_for_ai = [{"role": "system", "content": "You are a precise grader outputting structured JSON."}, {"role": "user", "content": [{"type": "text", "text": main_prompt_text}] + image_messages_for_ai}]
    try:
        response = _call_openai_with_retries(chat_client, messages_for_ai, os.getenv("AZURE_OPENAI_CHAT_DEPLOYMENT_NAME"), 0.4, 4000, {"type": "json_object"})
        parsed_result = json.loads(response.choices[0].message.content)
        overall_result = {"total_score": parsed_result.get("overall_total_score", "N/A"), "overall_feedback": parsed_result.get("overall_feedback", "N/A")}
        grading_breakdown = parsed_result.get("grades", [])
        return None, grading_breakdown, overall_result
    except Exception as e:
        print(f"An error occurred during AI grading: {e}"); import traceback; traceback.print_exc()
        return f"AI grading error: {e}", [], {"total_score": "N/A", "overall_feedback": f"AI grading failed: {e}"}
 
def generate_styled_excel_report(original_rubric_dataframe, grading_breakdown, overall_result):
    """
    Generates a styled Excel report with correct coloring, now fully dynamic.
    """
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        workbook, worksheet = writer.book, writer.book.add_worksheet('Analysis Results')
 
        col_map = getattr(original_rubric_dataframe, '_identified_columns', {})
        actual_criteria_col_name = col_map.get('criterion_col', 'Criterion')
        actual_category_col_name = col_map.get('category_col') # Might be None
        actual_max_score_col_name = col_map.get('max_score_col', 'Score')
 
        df_rubric_body = original_rubric_dataframe[~original_rubric_dataframe['is_summary_row']].copy()
        df_rubric_body.dropna(subset=[actual_criteria_col_name], inplace=True)
        report_columns = [col for col in original_rubric_dataframe.columns if col not in ['is_summary_row', 'index']]
        if 'AI Score' not in report_columns: report_columns.append('AI Score')
        if 'AI Comments' not in report_columns: report_columns.append('AI Comments')
 
        ai_grades_map = {grade.get('criterion_id'): grade for grade in grading_breakdown}
        df_rubric_body['AI Score'] = df_rubric_body['index'].map(lambda idx: ai_grades_map.get(idx, {}).get('score_achieved', 'N/A'))
        df_rubric_body['AI Comments'] = df_rubric_body['index'].map(lambda idx: ai_grades_map.get(idx, {}).get('comments', 'No AI feedback.'))
 
        final_report_rows_list = []
       
        # Only group and create subtotals if a category column was identified
        if actual_category_col_name and actual_category_col_name in df_rubric_body.columns:
            grouped = df_rubric_body.groupby(actual_category_col_name, sort=False)
            for category_name, group_df in grouped:
                for _, row in group_df.iterrows():
                    final_report_rows_list.append(row.to_dict())
               
                category_achieved = safe_numeric_score(pd.to_numeric(group_df['AI Score'], errors='coerce').sum())
                category_max = safe_numeric_score(pd.to_numeric(group_df[actual_max_score_col_name], errors='coerce').sum())
                subtotal_row = {col: "" for col in report_columns}
                subtotal_row[actual_category_col_name] = f"{category_name} Total"
                if actual_max_score_col_name: subtotal_row[actual_max_score_col_name] = category_max
                subtotal_row['AI Score'] = category_achieved
                final_report_rows_list.append(subtotal_row)
        else:
            # If no category column, just add all rows without subtotals
            for _, row in df_rubric_body.iterrows():
                final_report_rows_list.append(row.to_dict())
 
        report_df = pd.DataFrame(final_report_rows_list, columns=report_columns).fillna('')
       
        overall_achieved = safe_numeric_score(pd.to_numeric(df_rubric_body['AI Score'], errors='coerce').sum())
        overall_max = safe_numeric_score(pd.to_numeric(df_rubric_body[actual_max_score_col_name], errors='coerce').sum()) if actual_max_score_col_name else 0
 
        feedback_row = {col: "" for col in report_columns}; feedback_row[actual_criteria_col_name] = "Overall Feedback"; feedback_row['AI Comments'] = overall_result.get('overall_feedback', 'N/A')
        total_row = {col: "" for col in report_columns}; total_row[actual_criteria_col_name] = f"TOTAL MARKS OUT OF {int(overall_max)}"; total_row['AI Score'] = overall_achieved
        report_df = pd.concat([report_df, pd.DataFrame([feedback_row, total_row])], ignore_index=True)
 
        # Define formats
        header_format = workbook.add_format({'bold': True, 'border': 1, 'align': 'center', 'valign': 'vcenter', 'bg_color': '#D9D9D9', 'text_wrap': True})
        border_format = workbook.add_format({'border': 1, 'text_wrap': True, 'valign': 'top'})
        blue_format = workbook.add_format({'bg_color': '#DDEBF7', 'border': 1, 'bold': True, 'text_wrap': True, 'valign': 'top'})
        yellow_format = workbook.add_format({'bg_color': '#FFFF00', 'border': 1, 'bold': True, 'text_wrap': True, 'valign': 'top'})
 
        # Write to Excel with coloring logic
        worksheet.write_row(0, 0, report_df.columns, header_format)
        for r_idx, row_data in report_df.iterrows():
            fmt = border_format
            cat_cell_val = str(row_data.get(actual_category_col_name, '')).lower().strip()
            criterion_cell_val = str(row_data.get(actual_criteria_col_name, '')).lower().strip()
           
            if "total marks out of" in criterion_cell_val: fmt = yellow_format
            elif cat_cell_val.endswith(" total"): fmt = blue_format
           
            worksheet.write_row(r_idx + 1, 0, row_data.values, fmt)
 
        # Adjust column widths
        for i, col in enumerate(report_df.columns):
            width = max(report_df[col].astype(str).map(len).max(), len(col)) + 3
            if col == actual_criteria_col_name: width = max(width, 50)
            if col == 'AI Comments': width = max(width, 60)
            worksheet.set_column(i, i, width)
 
    output.seek(0)
    return output.getvalue(), report_df
 